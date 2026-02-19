#!/usr/bin/env python3
"""
PPT 템플릿을 사용한 표지 텍스트 오버레이
- PPT 템플릿에 이미지 삽입
- 텍스트만 변경
- PNG로 내보내기
"""

import subprocess
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN
import shutil

PROJECT_ROOT = Path(__file__).parent.parent.parent
TEMPLATE_PATH = PROJECT_ROOT / "content" / "templates" / "text_guide.pptx"


def apply_template(cover_image_path: str, title_text: str, output_path: str):
    """
    PPT 템플릿에 이미지와 텍스트를 적용하고 PNG로 내보내기

    Args:
        cover_image_path: 표지 이미지 경로
        title_text: 제목 텍스트 (예: "RICE", "CUCUMBER")
        output_path: 출력 PNG 경로
    """
    # 템플릿 복사해서 작업
    temp_pptx = Path(output_path).with_suffix('.pptx')
    shutil.copy(TEMPLATE_PATH, temp_pptx)

    # PPT 열기
    prs = Presentation(str(temp_pptx))
    slide = prs.slides[0]

    # 슬라이드 크기 확인
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 기존 이미지 제거하고 새 이미지 추가 (맨 뒤로)
    # 먼저 기존 그림 shape 찾기
    shapes_to_remove = []
    text_shape = None

    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            shapes_to_remove.append(shape)
        elif shape.has_text_frame:
            # 텍스트가 있는 shape 찾기
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text.strip():
                    text_shape = shape
                    break

    # 기존 이미지 제거
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

    # 새 이미지 추가 (전체 슬라이드 크기로)
    left = Emu(0)
    top = Emu(0)
    pic = slide.shapes.add_picture(
        cover_image_path,
        left, top,
        width=slide_width,
        height=slide_height
    )

    # 이미지를 맨 뒤로 보내기
    spTree = slide.shapes._spTree
    pic_element = pic._element
    spTree.remove(pic_element)
    spTree.insert(2, pic_element)  # 배경 바로 위에

    # 텍스트 변경
    if text_shape:
        for paragraph in text_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.text = title_text

    # 저장
    prs.save(str(temp_pptx))

    # PNG로 내보내기 (LibreOffice 또는 Keynote 사용)
    export_to_png(temp_pptx, output_path)

    # 임시 파일 정리
    temp_pptx.unlink()

    print(f"✅ 완료: {output_path}")


def export_to_png(pptx_path: Path, png_path: str):
    """PPTX를 PNG로 내보내기"""
    # macOS: Keynote나 sips 사용 불가, LibreOffice 사용
    # 또는 pdf2image 사용

    try:
        # LibreOffice로 PDF 변환 후 PNG
        pdf_path = pptx_path.with_suffix('.pdf')

        subprocess.run([
            '/Applications/LibreOffice.app/Contents/MacOS/soffice',
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', str(pptx_path.parent),
            str(pptx_path)
        ], check=True, capture_output=True)

        # PDF to PNG (ImageMagick)
        subprocess.run([
            'convert',
            '-density', '300',
            str(pdf_path) + '[0]',
            '-resize', '1080x1080',
            png_path
        ], check=True, capture_output=True)

        pdf_path.unlink()

    except Exception as e:
        print(f"⚠️  LibreOffice/ImageMagick 변환 실패: {e}")
        print("대체 방법: Puppeteer로 변환")
        # Fallback: 수동 변환 안내
        print(f"📌 수동 작업 필요: {pptx_path} → PNG 내보내기")


if __name__ == "__main__":
    if len(sys.argv) < 4:
        print("Usage: python apply_ppt_template.py <cover_image> <title> <output_path>")
        sys.exit(1)

    apply_template(sys.argv[1], sys.argv[2], sys.argv[3])
