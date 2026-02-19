#!/usr/bin/env python3
"""
PPT 템플릿을 사용한 단일 표지 텍스트 오버레이
python-pptx로 PPTX 생성 → PowerPoint AppleScript로 PNG 내보내기
"""

import subprocess
import sys
import time
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import shutil

PROJECT_ROOT = Path(__file__).parent.parent.parent
TEMPLATE_PATH = PROJECT_ROOT / "content" / "templates" / "text_guide.pptx"


def create_cover_pptx(cover_image_path: str, title_text: str, output_pptx: str):
    """
    PPT 템플릿에 이미지와 텍스트를 적용
    """
    cover_path = Path(cover_image_path)
    if not cover_path.exists():
        print(f"❌ 이미지 파일 없음: {cover_image_path}")
        return False

    # 템플릿 복사
    shutil.copy(TEMPLATE_PATH, output_pptx)

    # PPT 열기
    prs = Presentation(output_pptx)
    slide = prs.slides[0]

    # 슬라이드 크기
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 기존 이미지 shape 제거
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shapes_to_remove.append(shape)

    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

    # 새 이미지 추가 (전체 슬라이드 크기)
    pic = slide.shapes.add_picture(
        str(cover_path),
        Emu(0), Emu(0),
        width=slide_width,
        height=slide_height
    )

    # 이미지를 맨 뒤로 보내기
    spTree = slide.shapes._spTree
    pic_element = pic._element
    spTree.remove(pic_element)
    spTree.insert(2, pic_element)  # 배경 바로 위에 삽입

    # 텍스트 shape 찾아서 텍스트 변경
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip():  # 텍스트가 있는 경우만
                        run.text = title_text

    # 저장
    prs.save(output_pptx)
    print(f"✅ PPTX 생성: {output_pptx}")
    return True


def export_pptx_to_png_via_powerpoint(pptx_path: str, png_output_dir: str):
    """
    PowerPoint AppleScript로 PNG 내보내기
    """
    pptx_path = str(Path(pptx_path).resolve())
    png_output_dir = str(Path(png_output_dir).resolve())

    applescript = f'''
    tell application "Microsoft PowerPoint"
        activate
        open POSIX file "{pptx_path}"
        delay 1

        set theDoc to active presentation

        -- PNG로 내보내기
        save theDoc in POSIX file "{png_output_dir}/slide" as save as PNG

        close theDoc saving no
    end tell
    '''

    try:
        result = subprocess.run(
            ['osascript', '-e', applescript],
            capture_output=True,
            text=True,
            timeout=30
        )

        if result.returncode == 0:
            print(f"✅ PNG 내보내기 완료: {png_output_dir}")
            return True
        else:
            print(f"⚠️ AppleScript 오류: {result.stderr}")
            return False

    except subprocess.TimeoutExpired:
        print("⚠️ PowerPoint 응답 시간 초과")
        return False
    except Exception as e:
        print(f"⚠️ 오류: {e}")
        return False


def apply_cover_overlay(cover_image: str, title: str, output_folder: str, filename_prefix: str):
    """
    전체 프로세스: PPTX 생성 → PNG 내보내기
    """
    output_folder = Path(output_folder)
    output_folder.mkdir(parents=True, exist_ok=True)

    # 임시 PPTX 경로
    temp_pptx = output_folder / f"_temp_{filename_prefix}.pptx"

    # 1. PPTX 생성
    if not create_cover_pptx(cover_image, title, str(temp_pptx)):
        return False

    # 2. PNG 내보내기
    success = export_pptx_to_png_via_powerpoint(str(temp_pptx), str(output_folder))

    # 3. 파일명 정리 (PowerPoint는 Slide1.png 형식으로 내보냄)
    if success:
        slide_png = output_folder / "Slide1.png"
        if slide_png.exists():
            final_png = output_folder / f"{filename_prefix}_00.png"
            slide_png.rename(final_png)
            print(f"✅ 최종 파일: {final_png}")

    # 4. 임시 파일 정리
    if temp_pptx.exists():
        temp_pptx.unlink()

    return success


if __name__ == "__main__":
    if len(sys.argv) < 5:
        print("Usage: python apply_single_cover_overlay.py <cover_image> <title> <output_folder> <prefix>")
        print("Example: python apply_single_cover_overlay.py cover.png RICE ./output rice")
        sys.exit(1)

    apply_cover_overlay(sys.argv[1], sys.argv[2], sys.argv[3], sys.argv[4])
