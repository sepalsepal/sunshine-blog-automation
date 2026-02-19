#!/usr/bin/env python3
"""
표지용 PPTX 파일 생성 (PNG 내보내기는 수동)
- python-pptx로 PPTX 파일 생성
- 사용자가 파일 열어서 PNG로 내보내기
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
import shutil

PROJECT_ROOT = Path(__file__).parent.parent.parent
TEMPLATE_PATH = PROJECT_ROOT / "content" / "templates" / "text_guide.pptx"


def create_cover_pptx(cover_image_path: str, title_text: str, output_pptx_path: str):
    """
    PPT 템플릿에 이미지와 텍스트를 적용하여 PPTX 생성
    """
    cover_path = Path(cover_image_path)
    if not cover_path.exists():
        print(f"❌ 이미지 파일 없음: {cover_image_path}")
        return False

    # 템플릿 복사
    shutil.copy(TEMPLATE_PATH, output_pptx_path)

    # PPT 열기
    prs = Presentation(output_pptx_path)
    slide = prs.slides[0]

    # 슬라이드 크기
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 기존 이미지 shape 제거
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shapes_to_remove.append(shape)

    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

    # 새 이미지 추가 (전체 슬라이드 크기)
    pic = slide.shapes.add_picture(
        str(cover_path),
        Emu(0), Emu(0),
        width=slide_width,
        height=slide_height
    )

    # 이미지를 맨 뒤로 보내기
    spTree = slide.shapes._spTree
    pic_element = pic._element
    spTree.remove(pic_element)
    spTree.insert(2, pic_element)

    # 텍스트 shape 찾아서 텍스트 변경
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip():
                        run.text = title_text

    # 저장
    prs.save(output_pptx_path)
    print(f"✅ PPTX 생성 완료: {output_pptx_path}")
    print(f"📌 박편집: 이 파일을 열고 '파일 > 내보내기 > PNG'로 저장하세요")
    return True


if __name__ == "__main__":
    if len(sys.argv) < 4:
        print("Usage: python create_cover_pptx.py <cover_image> <title> <output_pptx>")
        print("Example: python create_cover_pptx.py cover.png RICE rice_cover.pptx")
        sys.exit(1)

    create_cover_pptx(sys.argv[1], sys.argv[2], sys.argv[3])
