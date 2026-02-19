#!/usr/bin/env python3
"""
PPTX → PNG 자동 변환 스크립트
LibreOffice + pypdfium2 사용

## 비용 안내 (LibreOffice)
| 항목 | 비용 |
|------|------|
| 초기 비용 | 무료 (오픈소스) |
| 월간 비용 | 없음 |
| 사용량 비용 | 없음 |

사용법:
    # 단일 PPTX 변환 (모든 슬라이드)
    python pptx_to_png.py input.pptx output_folder/

    # 특정 슬라이드만 변환
    python pptx_to_png.py input.pptx output_folder/ --slides 0,1,2

    # 템플릿에 이미지+텍스트 적용 후 변환
    python pptx_to_png.py --template --image cover.png --text "RICE" --output rice_00.png
"""

import subprocess
import sys
import tempfile
import shutil
from pathlib import Path
from typing import List, Optional

# pypdfium2, python-pptx 임포트
try:
    import pypdfium2 as pdfium
    from pptx import Presentation
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError as e:
    print(f"필수 패키지 설치 필요: {e}")
    print("pip install pypdfium2 python-pptx")
    sys.exit(1)

# 경로 설정
PROJECT_ROOT = Path(__file__).parent.parent.parent
TEMPLATE_PATH = PROJECT_ROOT / "content" / "templates" / "text_guide.pptx"
LIBREOFFICE_PATH = "/Applications/LibreOffice.app/Contents/MacOS/soffice"

# 출력 크기
OUTPUT_SIZE = 1080


def check_libreoffice() -> bool:
    """LibreOffice 설치 확인"""
    if not Path(LIBREOFFICE_PATH).exists():
        print("❌ LibreOffice가 설치되어 있지 않습니다.")
        print("설치: brew install --cask libreoffice")
        print("\n## 비용 안내 (LibreOffice)")
        print("| 항목 | 비용 |")
        print("|------|------|")
        print("| 초기 비용 | 무료 |")
        print("| 월간 비용 | 없음 |")
        return False
    return True


def pptx_to_pdf(pptx_path: Path, output_dir: Path) -> Optional[Path]:
    """PPTX → PDF 변환 (LibreOffice 사용)"""
    if not check_libreoffice():
        return None

    output_dir.mkdir(parents=True, exist_ok=True)

    cmd = [
        LIBREOFFICE_PATH,
        "--headless",
        "--convert-to", "pdf",
        "--outdir", str(output_dir),
        str(pptx_path)
    ]

    try:
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
        if result.returncode != 0:
            print(f"❌ LibreOffice 변환 실패: {result.stderr}")
            return None

        pdf_path = output_dir / f"{pptx_path.stem}.pdf"
        if pdf_path.exists():
            return pdf_path
        else:
            print(f"❌ PDF 파일 생성 실패")
            return None

    except subprocess.TimeoutExpired:
        print("❌ LibreOffice 타임아웃 (60초)")
        return None
    except Exception as e:
        print(f"❌ 오류: {e}")
        return None


def pdf_to_png(pdf_path: Path, output_dir: Path,
               slides: Optional[List[int]] = None,
               prefix: str = "slide") -> List[Path]:
    """
    PDF → PNG 변환 (pypdfium2 사용, 고품질 1080x1080)

    고품질 렌더링: 3x 스케일(3240px) → LANCZOS 다운샘플링(1080px)
    PD 수동 PPT 내보내기와 동일한 품질 달성
    """
    from PIL import Image

    output_dir.mkdir(parents=True, exist_ok=True)
    output_files = []

    # 고품질 렌더링 설정 (3x oversample)
    HQ_RENDER_SIZE = 3240  # 1080 * 3

    try:
        pdf = pdfium.PdfDocument(str(pdf_path))
        total_pages = len(pdf)

        # 변환할 페이지 결정
        if slides is None:
            pages_to_convert = range(total_pages)
        else:
            pages_to_convert = [s for s in slides if 0 <= s < total_pages]

        for page_num in pages_to_convert:
            page = pdf[page_num]

            # 페이지 크기 확인
            width, height = page.get_size()

            # 3x 고품질 렌더링
            scale = HQ_RENDER_SIZE / max(width, height)
            bitmap = page.render(scale=scale)
            pil_image = bitmap.to_pil()

            # LANCZOS 다운샘플링으로 1080x1080 출력
            final_image = pil_image.resize((OUTPUT_SIZE, OUTPUT_SIZE), Image.LANCZOS)

            # 저장
            output_path = output_dir / f"{prefix}_{page_num:02d}.png"
            final_image.save(output_path, "PNG")
            output_files.append(output_path)

            file_size = output_path.stat().st_size
            print(f"✅ 슬라이드 {page_num}: {output_path.name} ({file_size:,} bytes)")

        pdf.close()
        return output_files

    except Exception as e:
        print(f"❌ PDF→PNG 변환 오류: {e}")
        return []


def convert_pptx_to_png(pptx_path: Path, output_dir: Path,
                        slides: Optional[List[int]] = None,
                        prefix: str = "slide",
                        cleanup: bool = True) -> List[Path]:
    """
    PPTX → PNG 전체 파이프라인

    Args:
        pptx_path: 입력 PPTX 파일 경로
        output_dir: 출력 폴더
        slides: 변환할 슬라이드 번호 리스트 (None이면 전체)
        prefix: 출력 파일 접두사
        cleanup: 임시 PDF 파일 삭제 여부

    Returns:
        생성된 PNG 파일 경로 리스트
    """
    print(f"📄 입력: {pptx_path}")
    print(f"📁 출력: {output_dir}")

    # 1. PPTX → PDF
    print("\n[1/2] PPTX → PDF 변환 중...")
    pdf_path = pptx_to_pdf(pptx_path, output_dir)
    if not pdf_path:
        return []
    print(f"✅ PDF 생성: {pdf_path.name}")

    # 2. PDF → PNG
    print("\n[2/2] PDF → PNG 변환 중...")
    png_files = pdf_to_png(pdf_path, output_dir, slides, prefix)

    # 3. 임시 파일 정리
    if cleanup and pdf_path.exists():
        pdf_path.unlink()
        print(f"🗑️ 임시 PDF 삭제")

    print(f"\n✅ 완료! {len(png_files)}개 PNG 생성")
    return png_files


def create_slide_from_template(
    template_path: Path,
    slide_index: int,
    background_image: Path,
    text_content: str,
    output_pptx: Path,
    subtitle_content: str = None
) -> bool:
    """
    템플릿의 특정 슬라이드에 이미지와 텍스트를 적용

    Args:
        template_path: PPT 템플릿 경로
        slide_index: 사용할 슬라이드 번호 (0-based)
        background_image: 배경 이미지 경로
        text_content: 제목 텍스트
        output_pptx: 출력 PPTX 경로
        subtitle_content: 부제목 텍스트 (선택)
    """
    if not background_image.exists():
        print(f"❌ 이미지 파일 없음: {background_image}")
        return False

    # 템플릿 복사
    shutil.copy(template_path, output_pptx)

    # PPT 열기
    prs = Presentation(str(output_pptx))

    if slide_index >= len(prs.slides):
        print(f"❌ 슬라이드 {slide_index}가 없습니다 (총 {len(prs.slides)}개)")
        return False

    slide = prs.slides[slide_index]
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 기존 이미지 shape 제거
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shapes_to_remove.append(shape)

    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

    # 새 이미지 추가 (전체 슬라이드 크기)
    pic = slide.shapes.add_picture(
        str(background_image),
        Emu(0), Emu(0),
        width=slide_width,
        height=slide_height
    )

    # 이미지를 맨 뒤로 보내기
    spTree = slide.shapes._spTree
    pic_element = pic._element
    spTree.remove(pic_element)
    spTree.insert(2, pic_element)

    # 텍스트 변경 (제목 + 부제목)
    # 텍스트 박스 목록 수집 (shape 단위)
    text_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            has_text = False
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip():
                        has_text = True
                        break
            if has_text:
                text_shapes.append(shape)

    # 첫 번째 텍스트 박스 = 제목 (모든 run을 합쳐서 변경)
    if text_shapes:
        shape = text_shapes[0]
        for paragraph in shape.text_frame.paragraphs:
            runs = list(paragraph.runs)
            if runs:
                # 첫 번째 run에 전체 텍스트, 나머지 run은 빈 텍스트
                runs[0].text = text_content
                for run in runs[1:]:
                    run.text = ""

    # 두 번째 텍스트 박스 = 부제목 (있으면)
    if len(text_shapes) > 1:
        shape = text_shapes[1]
        for paragraph in shape.text_frame.paragraphs:
            runs = list(paragraph.runs)
            if runs:
                runs[0].text = subtitle_content if subtitle_content else ""
                for run in runs[1:]:
                    run.text = ""

    # 저장 (슬라이드 삭제 대신 PDF 변환 시 특정 슬라이드만 선택)
    prs.save(str(output_pptx))
    print(f"✅ PPTX 생성: {output_pptx}")
    return True


def create_overlay_png(
    template_path: Path,
    slide_index: int,
    background_image: Path,
    text_content: str,
    output_png: Path,
    subtitle_content: str = None,
    cleanup: bool = True
) -> bool:
    """
    템플릿 기반 PNG 생성 (이미지 + 텍스트 오버레이)

    전체 파이프라인:
    1. 템플릿 슬라이드에 이미지/텍스트 적용
    2. PPTX → PDF → PNG 변환

    Args:
        template_path: PPT 템플릿 경로
        slide_index: 사용할 슬라이드 번호 (0: 표지, 1-9: 본문)
        background_image: 배경 이미지 경로
        text_content: 제목 텍스트
        output_png: 출력 PNG 경로
        subtitle_content: 부제목 텍스트 (선택)
        cleanup: 임시 파일 삭제 여부
    """
    output_png = Path(output_png)
    output_dir = output_png.parent
    output_dir.mkdir(parents=True, exist_ok=True)

    # 임시 PPTX 경로
    temp_pptx = output_dir / f"_temp_{output_png.stem}.pptx"

    try:
        # 1. 템플릿에 이미지/텍스트 적용
        print(f"\n[1/3] 템플릿 적용 중... (슬라이드 {slide_index})")
        if not create_slide_from_template(
            template_path, slide_index, background_image, text_content, temp_pptx,
            subtitle_content=subtitle_content
        ):
            return False

        # 2. PPTX → PDF
        print("\n[2/3] PPTX → PDF 변환 중...")
        pdf_path = pptx_to_pdf(temp_pptx, output_dir)
        if not pdf_path:
            return False

        # 3. PDF → PNG (해당 슬라이드만 변환)
        print("\n[3/3] PDF → PNG 변환 중...")
        png_files = pdf_to_png(pdf_path, output_dir, slides=[slide_index], prefix=output_png.stem)

        if png_files:
            # 파일명 정리 (slide_00.png → 원하는 이름)
            generated = png_files[0]
            if generated != output_png:
                if output_png.exists():
                    output_png.unlink()
                generated.rename(output_png)
            print(f"\n✅ 최종 파일: {output_png}")
            return True

        return False

    finally:
        # 임시 파일 정리
        if cleanup:
            if temp_pptx.exists():
                temp_pptx.unlink()
            pdf_temp = output_dir / f"{temp_pptx.stem}.pdf"
            if pdf_temp.exists():
                pdf_temp.unlink()


# ============================================================
# CLI 인터페이스
# ============================================================

def main():
    import argparse

    parser = argparse.ArgumentParser(
        description="PPTX → PNG 변환 (LibreOffice + pypdfium2)",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
예시:
  # 전체 슬라이드 변환
  python pptx_to_png.py presentation.pptx ./output/

  # 특정 슬라이드만 변환 (0, 1, 2번)
  python pptx_to_png.py presentation.pptx ./output/ --slides 0,1,2

  # 템플릿 모드: 이미지+텍스트 → PNG
  python pptx_to_png.py --template --slide 0 --image cover.png --text "RICE" --output rice_00.png

비용 안내 (LibreOffice):
  | 항목 | 비용 |
  |------|------|
  | 초기 비용 | 무료 |
  | 월간 비용 | 없음 |
        """
    )

    # 기본 모드 인자
    parser.add_argument("input", nargs="?", help="입력 PPTX 파일")
    parser.add_argument("output", nargs="?", help="출력 폴더")
    parser.add_argument("--slides", help="변환할 슬라이드 번호 (쉼표 구분, 예: 0,1,2)")
    parser.add_argument("--prefix", default="slide", help="출력 파일 접두사")

    # 템플릿 모드 인자
    parser.add_argument("--template", action="store_true", help="템플릿 모드 사용")
    parser.add_argument("--slide", type=int, default=0, help="사용할 템플릿 슬라이드 번호")
    parser.add_argument("--image", help="배경 이미지 경로")
    parser.add_argument("--text", help="제목 텍스트")
    parser.add_argument("--subtitle", help="부제목 텍스트 (선택)")
    parser.add_argument("--output-png", dest="output_png", help="출력 PNG 경로")

    args = parser.parse_args()

    # 템플릿 모드
    if args.template:
        if not all([args.image, args.text, args.output_png]):
            print("❌ 템플릿 모드 필수 인자: --image, --text, --output-png")
            sys.exit(1)

        success = create_overlay_png(
            template_path=TEMPLATE_PATH,
            slide_index=args.slide,
            background_image=Path(args.image),
            text_content=args.text,
            output_png=Path(args.output_png),
            subtitle_content=args.subtitle
        )
        sys.exit(0 if success else 1)

    # 기본 모드 (PPTX → PNG)
    if not args.input or not args.output:
        parser.print_help()
        sys.exit(1)

    slides = None
    if args.slides:
        slides = [int(s.strip()) for s in args.slides.split(",")]

    png_files = convert_pptx_to_png(
        pptx_path=Path(args.input),
        output_dir=Path(args.output),
        slides=slides,
        prefix=args.prefix
    )

    sys.exit(0 if png_files else 1)


if __name__ == "__main__":
    main()
