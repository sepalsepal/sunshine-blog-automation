# 🎨 PPT 텍스트 오버레이 자동화 스크립트
# AI Crew 결과물(JSON)을 PPT 템플릿에 적용
#
# 사용법: python ppt_overlay.py peach
# 입력: peach_content.json + templates/text_guide.pptx
# 출력: outputs/peach/peach_slides.pptx

import json
import sys
import os
from pptx import Presentation
from pptx.util import Pt
from pptx.util import Inches
from copy import deepcopy

# ============================================================
# 🔧 설정
# ============================================================

TEMPLATE_PATH = "content/templates/text_guide.pptx"
OUTPUT_DIR = "content/images"

# ============================================================
# 📂 파일 로드 함수
# ============================================================

def load_content(topic: str) -> dict:
    """AI Crew 결과물 JSON 로드"""
    json_path = f"{topic}_content.json"

    if not os.path.exists(json_path):
        print(f"❌ 파일 없음: {json_path}")
        print(f"   먼저 'python ai_crew.py {topic}' 실행하세요.")
        sys.exit(1)

    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    print(f"✅ JSON 로드 완료: {json_path}")
    return data


def get_slides_text(data: dict) -> list:
    """JSON에서 슬라이드 텍스트 추출"""
    # 김작가의 텍스트 (text 섹션) 사용
    if "text" in data and "slides" in data["text"]:
        return data["text"]["slides"]
    # 없으면 기획안 (plan 섹션) 사용
    elif "plan" in data and "slides" in data["plan"]:
        return data["plan"]["slides"]
    else:
        print("❌ JSON에서 slides 데이터를 찾을 수 없습니다.")
        sys.exit(1)


# ============================================================
# 🖼️ PPT 처리 함수
# ============================================================

def apply_text_to_ppt(template_path: str, slides_data: list, output_path: str):
    """PPT 템플릿에 텍스트 적용"""

    # 템플릿 로드
    if not os.path.exists(template_path):
        print(f"❌ 템플릿 없음: {template_path}")
        sys.exit(1)

    prs = Presentation(template_path)
    print(f"✅ 템플릿 로드: {template_path}")
    print(f"   슬라이드 수: {len(prs.slides)}장")

    # 4장만 사용 (기존 7장 템플릿에서)
    slides_to_use = min(len(slides_data), len(prs.slides), 4)

    for i in range(slides_to_use):
        slide = prs.slides[i]
        slide_data = slides_data[i]

        slide_num = slide_data.get("slide", i + 1)
        slide_type = slide_data.get("type", "")

        print(f"\n📝 슬라이드 {slide_num} ({slide_type}) 처리 중...")

        # 표지 슬라이드
        if slide_type == "cover":
            title = slide_data.get("title", "")
            apply_cover_text(slide, title)

        # 본문 슬라이드
        else:
            main_text = slide_data.get("main_text", "")
            sub_text = slide_data.get("sub_text", "")
            apply_body_text(slide, main_text, sub_text)

    # 5장 이상 슬라이드 삭제 (4장만 유지)
    while len(prs.slides) > 4:
        rId = prs.slides._sldIdLst[4].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[4]

    # 출력 디렉토리 생성
    output_dir = os.path.dirname(output_path)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir)

    # 저장
    prs.save(output_path)
    print(f"\n✅ PPT 저장 완료: {output_path}")

    return prs


def apply_cover_text(slide, title: str):
    """표지 슬라이드에 타이틀 적용"""
    for shape in slide.shapes:
        if shape.has_text_frame:
            # 텍스트 박스 찾기
            text_frame = shape.text_frame
            if len(text_frame.paragraphs) > 0:
                para = text_frame.paragraphs[0]
                if len(para.runs) > 0:
                    # 기존 텍스트 확인
                    old_text = para.runs[0].text
                    # 타이틀 교체
                    para.runs[0].text = title
                    print(f"   표지 타이틀: '{old_text}' → '{title}'")
                    return

    print(f"   ⚠️ 표지 텍스트 박스를 찾지 못했습니다.")


def apply_body_text(slide, main_text: str, sub_text: str):
    """본문 슬라이드에 메인/서브 텍스트 적용"""
    text_boxes = []

    # 모든 텍스트 박스 수집
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_boxes.append(shape)

    # 위치 기준 정렬 (위에서 아래로)
    text_boxes.sort(key=lambda s: s.top)

    # 텍스트 박스가 2개 이상이면 (메인 + 서브)
    main_applied = False
    sub_applied = False

    for shape in text_boxes:
        text_frame = shape.text_frame
        for para in text_frame.paragraphs:
            for run in para.runs:
                old_text = run.text.strip()

                # 빈 텍스트 스킵
                if not old_text:
                    continue

                # 메인 텍스트 (보통 더 큰 폰트)
                if not main_applied and run.font.size and run.font.size >= Pt(30):
                    run.text = main_text
                    print(f"   메인: '{old_text}' → '{main_text}'")
                    main_applied = True

                # 서브 텍스트 (보통 더 작은 폰트)
                elif not sub_applied and run.font.size and run.font.size < Pt(30):
                    run.text = sub_text
                    print(f"   서브: '{old_text}' → '{sub_text}'")
                    sub_applied = True

    if not main_applied:
        print(f"   ⚠️ 메인 텍스트 박스를 찾지 못했습니다.")
    if not sub_applied:
        print(f"   ⚠️ 서브 텍스트 박스를 찾지 못했습니다.")


# ============================================================
# 🖼️ PNG 내보내기 (LibreOffice 사용)
# ============================================================

def export_to_png(pptx_path: str, output_dir: str):
    """PPT를 PNG 이미지로 내보내기"""
    import subprocess

    print(f"\n🖼️ PNG 내보내기 중...")

    # LibreOffice로 PDF 변환
    try:
        subprocess.run([
            "soffice", "--headless", "--convert-to", "pdf",
            "--outdir", output_dir, pptx_path
        ], check=True, capture_output=True)

        pdf_path = pptx_path.replace(".pptx", ".pdf")
        pdf_path = os.path.join(output_dir, os.path.basename(pdf_path))

        # PDF를 PNG로 변환
        subprocess.run([
            "pdftoppm", "-png", "-r", "300", pdf_path,
            os.path.join(output_dir, "slide")
        ], check=True, capture_output=True)

        print(f"✅ PNG 내보내기 완료: {output_dir}/slide-*.png")

    except FileNotFoundError:
        print("⚠️ LibreOffice 또는 pdftoppm이 설치되어 있지 않습니다.")
        print("   수동으로 PPT를 열어 PNG로 내보내세요.")
    except subprocess.CalledProcessError as e:
        print(f"⚠️ PNG 변환 중 오류: {e}")


# ============================================================
# 🚀 메인 함수
# ============================================================

def main():
    if len(sys.argv) < 2:
        print("사용법: python ppt_overlay.py <주제>")
        print("예시: python ppt_overlay.py peach")
        sys.exit(1)

    topic = sys.argv[1]

    print("=" * 60)
    print(f"🎨 PPT 텍스트 오버레이 자동화")
    print(f"📝 주제: {topic}")
    print("=" * 60)

    # 1. JSON 로드
    data = load_content(topic)
    slides_data = get_slides_text(data)

    print(f"\n📋 슬라이드 데이터:")
    for s in slides_data[:4]:  # 4장만 표시
        slide_num = s.get("slide", "?")
        slide_type = s.get("type", "?")
        title = s.get("title", "")
        main = s.get("main_text", "")[:20]
        print(f"   [{slide_num}] {slide_type}: {title or main}")

    # 2. PPT 적용
    output_dir = os.path.join(OUTPUT_DIR, topic)
    output_path = os.path.join(output_dir, f"{topic}_slides.pptx")

    apply_text_to_ppt(TEMPLATE_PATH, slides_data, output_path)

    # 3. PNG 내보내기 (선택)
    if len(sys.argv) > 2 and sys.argv[2] == "--png":
        export_to_png(output_path, output_dir)

    # 결과 요약
    print("\n" + "=" * 60)
    print("✅ 완료!")
    print(f"📁 PPT 파일: {output_path}")
    print(f"📁 PNG 변환: python ppt_overlay.py {topic} --png")
    print("=" * 60)


if __name__ == "__main__":
    main()
