#!/usr/bin/env python3
"""
PPT 템플릿 기반 텍스트 오버레이 스크립트
- content/templates/text_guide.pptx 템플릿 사용
- LibreOffice로 PNG 변환 (1080x1080)

담당: 박편집
검수: 김감독
"""

import json
import os
import shutil
import subprocess
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

ROOT = Path(__file__).parent.parent.parent
TEMPLATE_PATH = ROOT / "content/templates/text_guide.pptx"
TEMP_DIR = ROOT / "content/templates/temp"


def create_slide_with_template(
    bg_image_path: str,
    title: str,
    subtitle: str,
    slide_type: str,
    output_path: str
):
    """
    PPT 템플릿으로 슬라이드 생성 후 PNG 내보내기

    Args:
        bg_image_path: 배경 이미지 경로
        title: 제목 텍스트
        subtitle: 부제목 텍스트
        slide_type: cover, danger, content_bottom, cta
        output_path: 출력 PNG 경로
    """
    print(f"📝 박편집입니다. PPT 템플릿으로 작업합니다.")
    print(f"   배경: {Path(bg_image_path).name}")
    print(f"   제목: {title}")
    print(f"   타입: {slide_type}")

    # 템플릿 열기
    prs = Presentation(str(TEMPLATE_PATH))
    slide = prs.slides[0]

    # 슬라이드 크기 확인 (18.06cm x 18.06cm = 1080x1080 at 150dpi 근사)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 1. 배경 이미지 삽입 (맨 뒤로)
    # 기존 이미지 shape 찾아서 교체 또는 새로 추가
    bg_shape = None
    for shape in slide.shapes:
        if shape.shape_type == 13:  # Picture
            bg_shape = shape
            break

    if bg_shape:
        # 기존 배경 이미지 삭제
        sp = bg_shape._element
        sp.getparent().remove(sp)

    # 새 배경 이미지 추가 (전체 슬라이드 크기)
    slide.shapes.add_picture(
        bg_image_path,
        Emu(0), Emu(0),
        width=slide_width,
        height=slide_height
    )

    # 배경 이미지를 맨 뒤로 보내기
    bg_pic = slide.shapes[-1]
    spTree = slide.shapes._spTree
    sp = bg_pic._element
    spTree.remove(sp)
    spTree.insert(2, sp)  # 맨 앞에 삽입 (z-order 맨 뒤)

    # 2. 텍스트 shape 찾아서 업데이트
    title_shape = None
    subtitle_shape = None

    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            # 기존 템플릿의 제목/부제목 찾기
            if "TITLE" in text.upper() or "제목" in text or len(text) < 20:
                if not title_shape:
                    title_shape = shape
                elif not subtitle_shape:
                    subtitle_shape = shape

    # 3. 제목 텍스트 설정
    if title_shape:
        tf = title_shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title

        # 스타일 설정 (타입별)
        font = run.font
        font.name = "Arial"
        font.bold = True

        if slide_type == "cover":
            font.size = Pt(54)
            font.color.rgb = RGBColor(255, 255, 255)
        elif slide_type == "danger":
            font.size = Pt(48)
            font.color.rgb = RGBColor(255, 107, 107)  # #FF6B6B 코랄 레드
        elif slide_type == "cta":
            font.size = Pt(48)
            font.color.rgb = RGBColor(255, 217, 61)  # #FFD93D 노란색
        else:
            font.size = Pt(48)
            font.color.rgb = RGBColor(255, 255, 255)

        p.alignment = PP_ALIGN.CENTER

    # 4. 부제목 텍스트 설정
    if subtitle_shape and subtitle:
        tf = subtitle_shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = subtitle

        font = run.font
        font.name = "Arial"
        font.size = Pt(24)
        font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    # 5. 임시 PPTX 저장
    TEMP_DIR.mkdir(parents=True, exist_ok=True)
    temp_pptx = TEMP_DIR / "temp_slide.pptx"
    prs.save(str(temp_pptx))

    # 6. LibreOffice로 PNG 변환
    print(f"   PNG 변환 중...")

    # soffice 명령어로 PDF 변환 후 PNG로
    subprocess.run([
        "/opt/homebrew/bin/soffice",
        "--headless",
        "--convert-to", "png",
        "--outdir", str(TEMP_DIR),
        str(temp_pptx)
    ], capture_output=True)

    # 변환된 PNG 찾기
    temp_png = TEMP_DIR / "temp_slide.png"

    if temp_png.exists():
        # 1080x1080으로 리사이즈
        img = Image.open(temp_png)
        img = img.resize((1080, 1080), Image.LANCZOS)

        # 출력 경로로 저장
        output_dir = Path(output_path).parent
        output_dir.mkdir(parents=True, exist_ok=True)
        img.save(output_path, quality=95)

        print(f"   ✅ 저장: {Path(output_path).name}")

        # 임시 파일 정리
        temp_pptx.unlink(missing_ok=True)
        temp_png.unlink(missing_ok=True)

        return True
    else:
        print(f"   ❌ PNG 변환 실패")
        return False


def render_grape_content():
    """포도 콘텐츠 PPT 템플릿으로 렌더링"""

    print("="*60)
    print("📝 박편집입니다. 포도 콘텐츠 PPT 템플릿 작업 시작합니다.")
    print("="*60)

    grape_dir = ROOT / "content/images/025_grape_포도"
    output_dir = grape_dir / "_final"
    output_dir.mkdir(parents=True, exist_ok=True)

    # 텍스트 설정 로드
    text_config = json.loads(
        (ROOT / "config/settings/grape_text.json").read_text(encoding='utf-8')
    )

    for slide in text_config:
        slide_num = slide["slide"]
        slide_type = slide["type"]
        title = slide["title"]
        subtitle = slide.get("subtitle", "")

        # 소스 이미지
        src_image = grape_dir / f"grape_0{slide_num}.png"
        output_path = output_dir / f"grape_0{slide_num}.png"

        if not src_image.exists():
            print(f"⚠️ 소스 이미지 없음: {src_image}")
            continue

        print(f"\n📌 Slide {slide_num} [{slide_type}]")

        success = create_slide_with_template(
            bg_image_path=str(src_image),
            title=title,
            subtitle=subtitle,
            slide_type=slide_type,
            output_path=str(output_path)
        )

        if not success:
            print(f"   ❌ 슬라이드 {slide_num} 실패")

    print("\n" + "="*60)
    print("📝 박편집입니다. 작업 완료. 김감독님 검수 부탁드립니다.")
    print("="*60)


if __name__ == "__main__":
    render_grape_content()
